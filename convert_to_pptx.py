#!/usr/bin/env python3
"""
Convert ADA Competitive Analysis to PowerPoint Presentation
Creates a professional PPTX with charts, data, and analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from datetime import datetime
import os

def create_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "ADA Competitive Analysis"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "vs FICO Falcon & SAS Fraud Management\n" + datetime.now().strftime("%B %Y")
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def create_executive_summary_slide(prs):
    """Create executive summary slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    content.text = """• ADA leads in ALL 4 key metrics:
  - Accuracy: 92% (vs 85% FICO, 80% SAS)
  - Speed: 95% (vs 75% FICO, 85% SAS)
  - Explainability: 90% (vs 60% FICO, 70% SAS)
  - Cost Efficiency: 80% (vs 65% FICO, 70% SAS)

• Superior real-time processing (< 1 second vs 1-5 seconds)
• Best explainability for regulatory compliance
• Lower total cost of ownership
• Lightweight infrastructure requirements"""
    
    # Format text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        if paragraph.text.startswith("• ADA leads"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(31, 119, 180)
    
    return slide

def create_performance_comparison_slide(prs):
    """Create performance comparison slide with chart"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Performance Comparison"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create bar chart
    metrics = ['Accuracy', 'Speed', 'Explainability', 'Cost Efficiency']
    ada_values = [92, 95, 90, 80]
    fico_values = [85, 75, 60, 65]
    sas_values = [80, 85, 70, 70]
    
    x = np.arange(len(metrics))
    width = 0.25
    
    fig, ax = plt.subplots(figsize=(10, 6))
    bars1 = ax.bar(x - width, ada_values, width, label='ADA', color='#4facfe')
    bars2 = ax.bar(x, fico_values, width, label='FICO Falcon', color='#f093fb')
    bars3 = ax.bar(x + width, sas_values, width, label='SAS Fraud', color='#f5576c')
    
    ax.set_xlabel('Metrics')
    ax.set_ylabel('Score (%)')
    ax.set_title('Performance Comparison Across Key Metrics')
    ax.set_xticks(x)
    ax.set_xticklabels(metrics)
    ax.legend()
    ax.set_ylim(0, 100)
    
    # Add value labels on bars
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 1,
                   f'{int(height)}%', ha='center', va='bottom')
    
    plt.tight_layout()
    plt.savefig('performance_chart.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add chart to slide
    slide.shapes.add_picture('performance_chart.png', Inches(1), Inches(2), Inches(8), Inches(5))
    
    # Clean up
    os.remove('performance_chart.png')
    
    return slide

def create_radar_chart_slide(prs):
    """Create radar chart slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Radar Chart Analysis"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create radar chart
    categories = ['Accuracy', 'Speed', 'Explainability', 'Cost Efficiency']
    ada_values = [92, 95, 90, 80]
    fico_values = [85, 75, 60, 65]
    sas_values = [80, 85, 70, 70]
    
    # Number of variables
    N = len(categories)
    
    # What will be the angle of each axis in the plot
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]
    
    # Initialize the spider plot
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
    
    # Plot data
    ada_values += ada_values[:1]
    fico_values += fico_values[:1]
    sas_values += sas_values[:1]
    
    ax.plot(angles, ada_values, 'o-', linewidth=2, label='ADA', color='#4facfe')
    ax.fill(angles, ada_values, alpha=0.25, color='#4facfe')
    
    ax.plot(angles, fico_values, 'o-', linewidth=2, label='FICO Falcon', color='#f093fb')
    ax.fill(angles, fico_values, alpha=0.25, color='#f093fb')
    
    ax.plot(angles, sas_values, 'o-', linewidth=2, label='SAS Fraud', color='#f5576c')
    ax.fill(angles, sas_values, alpha=0.25, color='#f5576c')
    
    # Fix axis to go in the right order and start at 12 o'clock
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    
    # Draw axis lines for each angle and label
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories)
    
    # Add legend
    ax.legend(loc='upper right', bbox_to_anchor=(0.1, 0.1))
    
    plt.tight_layout()
    plt.savefig('radar_chart.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add chart to slide
    slide.shapes.add_picture('radar_chart.png', Inches(1), Inches(2), Inches(8), Inches(6))
    
    # Clean up
    os.remove('radar_chart.png')
    
    return slide

def create_feature_comparison_slide(prs):
    """Create feature comparison slide with table"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Feature Comparison"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Create table data
    data = {
        'Feature': [
            'Detection Accuracy',
            'Processing Speed',
            'Real-time Processing',
            'Explainability',
            'Cost Efficiency',
            'Infrastructure Requirements',
            'Integration Complexity',
            'Customization',
            'Scalability',
            'Support & Maintenance'
        ],
        'ADA': [
            '92%',
            '95%',
            '< 1 second',
            '90%',
            '80%',
            'Lightweight',
            'Low',
            'High',
            'Excellent',
            'Low'
        ],
        'FICO Falcon': [
            '85%',
            '75%',
            '3-5 seconds',
            '60%',
            '65%',
            'Heavy',
            'High',
            'Moderate',
            'Good',
            'High'
        ],
        'SAS Fraud': [
            '80%',
            '85%',
            '1-2 seconds',
            '70%',
            '70%',
            'Moderate',
            'Medium',
            'High',
            'Good',
            'Medium'
        ]
    }
    
    df = pd.DataFrame(data)
    
    # Create table
    rows, cols = len(df) + 1, len(df.columns)
    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(5)).table
    
    # Add headers
    for i, col in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
    
    # Add data
    for i, row in df.iterrows():
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            
            # Highlight ADA advantages
            if j > 0 and str(value) in ['92%', '95%', '< 1 second', '90%', '80%', 'Lightweight', 'Low', 'High', 'Excellent']:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green
    
    return slide

def create_cost_analysis_slide(prs):
    """Create cost analysis slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Cost Analysis"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create cost comparison chart
    cost_components = ['Licensing', 'Infrastructure', 'Integration', 'Maintenance', 'Training']
    ada_costs = [100, 80, 60, 70, 50]
    fico_costs = [150, 120, 100, 90, 80]
    sas_costs = [130, 100, 80, 85, 70]
    
    x = np.arange(len(cost_components))
    width = 0.25
    
    fig, ax = plt.subplots(figsize=(10, 6))
    bars1 = ax.bar(x - width, ada_costs, width, label='ADA', color='#4facfe')
    bars2 = ax.bar(x, fico_costs, width, label='FICO Falcon', color='#f093fb')
    bars3 = ax.bar(x + width, sas_costs, width, label='SAS Fraud', color='#f5576c')
    
    ax.set_xlabel('Cost Components')
    ax.set_ylabel('Relative Cost Index')
    ax.set_title('Cost Comparison (Relative Index)')
    ax.set_xticks(x)
    ax.set_xticklabels(cost_components)
    ax.legend()
    
    # Add value labels on bars
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 2,
                   f'{int(height)}', ha='center', va='bottom')
    
    plt.tight_layout()
    plt.savefig('cost_chart.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add chart to slide
    slide.shapes.add_picture('cost_chart.png', Inches(1), Inches(2), Inches(8), Inches(5))
    
    # Clean up
    os.remove('cost_chart.png')
    
    return slide

def create_competitive_advantages_slide(prs):
    """Create competitive advantages slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "ADA Competitive Advantages"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    content.text = """⚡ Speed & Performance:
• < 1 second processing time (vs 1-5 seconds)
• 500+ TPS throughput
• Real-time fraud detection
• Low latency architecture

🧠 AI & Explainability:
• 90% explainability score
• XAI integration
• Regulatory compliance
• Transparent decisions

💰 Cost Efficiency:
• Lower TCO than competitors
• Lightweight infrastructure
• Easy integration
• Reduced maintenance

🎯 Accuracy & Reliability:
• 92% accuracy rate
• Advanced ML models
• Continuous learning
• Adaptive algorithms"""
    
    # Format text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        if paragraph.text.startswith("⚡") or paragraph.text.startswith("🧠") or paragraph.text.startswith("💰") or paragraph.text.startswith("🎯"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(31, 119, 180)
    
    return slide

def create_market_positioning_slide(prs):
    """Create market positioning slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Market Positioning"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create scatter plot
    fig, ax = plt.subplots(figsize=(8, 6))
    
    # Data points
    accuracy = [92, 85, 80]
    speed = [95, 75, 85]
    labels = ['ADA', 'FICO', 'SAS']
    colors = ['#4facfe', '#f093fb', '#f5576c']
    sizes = [300, 200, 200]
    
    scatter = ax.scatter(accuracy, speed, c=colors, s=sizes, alpha=0.7)
    
    # Add labels
    for i, label in enumerate(labels):
        ax.annotate(label, (accuracy[i], speed[i]), xytext=(5, 5), textcoords='offset points')
    
    ax.set_xlabel('Accuracy (%)')
    ax.set_ylabel('Speed (%)')
    ax.set_title('Market Positioning: Accuracy vs Speed')
    ax.grid(True, alpha=0.3)
    ax.set_xlim(75, 100)
    ax.set_ylim(70, 100)
    
    plt.tight_layout()
    plt.savefig('positioning_chart.png', dpi=300, bbox_inches='tight')
    plt.close()
    
    # Add chart to slide
    slide.shapes.add_picture('positioning_chart.png', Inches(1), Inches(2), Inches(8), Inches(5))
    
    # Clean up
    os.remove('positioning_chart.png')
    
    return slide

def create_recommendations_slide(prs):
    """Create recommendations slide"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title = slide.shapes.title
    title.text = "Strategic Recommendations"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content
    content = slide.placeholders[1]
    content.text = """🚀 Go-to-Market Strategy:
1. Target High-Value Segments:
   • Financial institutions requiring real-time processing
   • E-commerce platforms with high transaction volumes
   • Regulated industries needing explainability

2. Competitive Messaging:
   • Emphasize speed advantage (5x faster than FICO)
   • Highlight cost savings (33% vs FICO, 23% vs SAS)
   • Showcase explainability for compliance

📈 Growth Opportunities:
3. Market Expansion:
   • Target mid-market companies priced out by FICO/SAS
   • Focus on emerging markets with cost sensitivity
   • Leverage cloud-native architecture advantage

4. Product Development:
   • Continue improving accuracy beyond 92%
   • Enhance explainability features
   • Develop industry-specific solutions

🎯 Key Recommendation:
Position ADA as the 'Next-Generation Fraud Detection Platform' with superior performance across all metrics, modern cloud-native architecture, lower TCO, and better regulatory compliance."""
    
    # Format text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        if paragraph.text.startswith("🚀") or paragraph.text.startswith("📈") or paragraph.text.startswith("🎯"):
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(31, 119, 180)
    
    return slide

def main():
    """Create the complete PowerPoint presentation"""
    print("📊 Creating ADA Competitive Analysis PowerPoint...")
    
    # Create presentation
    prs = Presentation()
    
    # Add slides
    print("📝 Adding slides...")
    create_title_slide(prs)
    create_executive_summary_slide(prs)
    create_performance_comparison_slide(prs)
    create_radar_chart_slide(prs)
    create_feature_comparison_slide(prs)
    create_cost_analysis_slide(prs)
    create_competitive_advantages_slide(prs)
    create_market_positioning_slide(prs)
    create_recommendations_slide(prs)
    
    # Save presentation
    filename = f"ADA_Competitive_Analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)
    
    print(f"✅ PowerPoint created successfully: {filename}")
    print(f"📁 File location: {os.path.abspath(filename)}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return filename

if __name__ == "__main__":
    main() 